from __future__ import annotations

import subprocess
import shutil
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader

from .models import ExtractedDocument

DEEP_EXTENSIONS = {
    ".jpg",
    ".jpeg",
    ".png",
    ".pdf",
    ".ppt",
    ".pptx",
    ".xls",
    ".xlsx",
    ".doc",
    ".docx",
}
MODERN_OFFICE = {".pptx", ".xlsx", ".docx"}
OFFICE_EXTENSIONS = {".ppt", ".pptx", ".xls", ".xlsx", ".doc", ".docx"}


def _run(command: list[str]) -> None:
    subprocess.run(command, check=True, capture_output=True, text=True)


def _pdf_text(path: Path) -> tuple[str, int]:
    reader = PdfReader(str(path))
    text = "\n\n".join(page.extract_text() or "" for page in reader.pages)
    return text, len(reader.pages)


def _modern_office_text(path: Path) -> str:
    if path.suffix.lower() == ".docx":
        doc = Document(path)
        return "\n".join(paragraph.text for paragraph in doc.paragraphs)
    if path.suffix.lower() == ".pptx":
        deck = Presentation(path)
        return "\n".join(
            shape.text
            for slide in deck.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
    workbook = load_workbook(path, read_only=True, data_only=True)
    chunks: list[str] = []
    for sheet in workbook.worksheets:
        chunks.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            values = [str(value) for value in row if value is not None]
            if values:
                chunks.append("\t".join(values))
    return "\n".join(chunks)


def _render_pdf(pdf_path: Path, output_dir: Path) -> tuple[list[str], str | None]:
    prefix = output_dir / "page"
    _run(
        [
            "pdftoppm",
            "-jpeg",
            "-f",
            "1",
            "-l",
            "8",
            "-r",
            "120",
            str(pdf_path),
            str(prefix),
        ]
    )
    images = sorted(str(item) for item in output_dir.glob("page-*.jpg"))
    return images, images[0] if images else None


def extract_document(
    source: Path,
    work_dir: Path,
    max_bytes: int,
    max_pages: int,
) -> ExtractedDocument:
    suffix = source.suffix.lower()
    if suffix not in DEEP_EXTENSIONS:
        return ExtractedDocument(
            text=source.name,
            deep_analysis_eligible=False,
            reason_skipped="unsupported_metadata_only",
        )
    if source.stat().st_size > max_bytes:
        return ExtractedDocument(
            text=source.name,
            deep_analysis_eligible=False,
            reason_skipped="over_100mb",
        )

    if suffix in {".jpg", ".jpeg", ".png"}:
        image = Image.open(source)
        thumbnail = work_dir / "thumbnail.jpg"
        preview_image = image.convert("RGB")
        preview_image.thumbnail((1280, 1280))
        preview_image.save(thumbnail, "JPEG", quality=88)
        return ExtractedDocument(
            text=source.name,
            page_count=1,
            thumbnail=str(thumbnail),
            image_paths=[str(source)],
        )

    preview: Path | None = source if suffix == ".pdf" else None
    if suffix in OFFICE_EXTENSIONS:
        office_binary = shutil.which("libreoffice") or shutil.which("soffice")
        if not office_binary:
            raise RuntimeError("LibreOffice executable not found")
        _run(
            [
                office_binary,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(work_dir),
                str(source),
            ]
        )
        preview = work_dir / f"{source.stem}.pdf"

    text = ""
    page_count: int | None = None
    if suffix in MODERN_OFFICE:
        text = _modern_office_text(source)
    if preview and preview.exists():
        pdf_text, page_count = _pdf_text(preview)
        text = text or pdf_text
    if page_count is not None and page_count > max_pages:
        return ExtractedDocument(
            text=source.name,
            page_count=page_count,
            preview_pdf=str(preview) if preview else None,
            deep_analysis_eligible=False,
            reason_skipped="over_100_pages",
        )

    images: list[str] = []
    thumbnail: str | None = None
    if preview and preview.exists():
        images, thumbnail = _render_pdf(preview, work_dir)
    return ExtractedDocument(
        text=text[:300_000],
        page_count=page_count,
        preview_pdf=str(preview) if preview and suffix != ".pdf" else None,
        thumbnail=thumbnail,
        image_paths=images,
    )
